#!/usr/bin/env python3
"""
Build a PPTX by placing each rendered slide image onto a blank slide.
Also write a run manifest with content and validation summaries.
"""

from __future__ import annotations

import hashlib
import json
import sys
from collections import Counter
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

BASE_PATH = Path(__file__).resolve().parent.parent
OUTPUT_DIR = BASE_PATH / "output"
IMAGES_DIR = OUTPUT_DIR / "images"
HTML_DIR = OUTPUT_DIR / "slides"
VALIDATION_DIR = OUTPUT_DIR / "validation"
PPT_OUTPUT = OUTPUT_DIR / "presentation.pptx"
MANIFEST_OUTPUT = OUTPUT_DIR / "manifest.json"
SPEC_CANDIDATES = [
    BASE_PATH / "deck-spec.json",
    BASE_PATH / "presentation-content.json",
    BASE_PATH / "content" / "deck-spec.json",
    BASE_PATH / "content" / "presentation-content.json",
]
VALIDATION_REPORTS = {
    "presentation_content": VALIDATION_DIR / "presentation-content-report.json",
    "deck_quality": VALIDATION_DIR / "deck-quality-report.json",
}

# PPT size: 16:9 (10in x 5.625in)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def load_json(path: Path):
    with path.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def find_spec_path() -> Path | None:
    for candidate in SPEC_CANDIDATES:
        if candidate.exists():
            return candidate
    return None


def get_image_files() -> list[Path]:
    if not IMAGES_DIR.exists():
        print(f"Image directory not found: {IMAGES_DIR}")
        print("Run: python src/render.py")
        sys.exit(1)

    image_files = sorted(IMAGES_DIR.glob("*.png"))
    if not image_files:
        print("No PNG files found. Run: python src/render.py")
        sys.exit(1)

    return image_files


def create_ppt(image_files: list[Path]) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]

    for image_file in image_files:
        slide = prs.slides.add_slide(blank_layout)
        try:
            slide.shapes.add_picture(
                str(image_file),
                Inches(0),
                Inches(0),
                width=SLIDE_WIDTH,
                height=SLIDE_HEIGHT,
            )
            print(f"  Added: {image_file.name}")
        except Exception as exc:
            print(f"  Failed: {image_file.name} - {exc}")
            sys.exit(1)

    return prs


def build_manifest(image_files: list[Path]) -> dict[str, object]:
    content_path = find_spec_path()
    content = None
    if content_path is not None:
        try:
            content = load_json(content_path)
        except Exception:
            content = None

    content_summary: dict[str, object] = {}
    if isinstance(content, dict):
        slides = content.get("slides", [])
        layout_counts = Counter()
        theme_counts = Counter()
        if isinstance(slides, list):
            for slide in slides:
                if not isinstance(slide, dict):
                    continue
                layout_id = slide.get("layout_id")
                theme_id = slide.get("theme")
                if isinstance(layout_id, str):
                    layout_counts[layout_id] += 1
                if isinstance(theme_id, str):
                    theme_counts[theme_id] += 1

        content_summary = {
            "path": str(content_path),
            "sha256": file_sha256(content_path),
            "version": content.get("version"),
            "source": content.get("source"),
            "slide_count": len(slides) if isinstance(slides, list) else 0,
            "layout_counts": dict(layout_counts),
            "theme_counts": dict(theme_counts),
        }

    html_files = sorted(HTML_DIR.glob("*.html")) if HTML_DIR.exists() else []
    validation_reports: dict[str, object] = {}
    for name, report_path in VALIDATION_REPORTS.items():
        if report_path.exists():
            try:
                validation_reports[name] = load_json(report_path)
            except Exception as exc:
                validation_reports[name] = {
                    "status": "error",
                    "path": str(report_path),
                    "error": str(exc),
                }
        else:
            validation_reports[name] = {
                "status": "missing",
                "path": str(report_path),
            }

    manifest: dict[str, object] = {
        "version": 1,
        "generated_at": datetime.now().astimezone().isoformat(timespec="seconds"),
        "content": content_summary,
        "artifacts": {
            "slides_html": {
                "path": str(HTML_DIR),
                "count": len(html_files),
            },
            "images": {
                "path": str(IMAGES_DIR),
                "count": len(image_files),
                "files": [
                    {
                        "name": image_file.name,
                        "size_bytes": image_file.stat().st_size,
                        "sha256": file_sha256(image_file),
                    }
                    for image_file in image_files
                ],
            },
            "pptx": {
                "path": str(PPT_OUTPUT),
                "exists": PPT_OUTPUT.exists(),
                "size_bytes": PPT_OUTPUT.stat().st_size if PPT_OUTPUT.exists() else 0,
                "sha256": file_sha256(PPT_OUTPUT) if PPT_OUTPUT.exists() else None,
            },
        },
        "validation": validation_reports,
    }

    contact_sheet_candidates = [
        OUTPUT_DIR / "contact-sheet.png",
        OUTPUT_DIR / "contact-sheet-2.png",
    ]
    contact_sheet = next((path for path in contact_sheet_candidates if path.exists()), None)
    if contact_sheet is not None:
        manifest["artifacts"]["contact_sheet"] = {
            "path": str(contact_sheet),
            "size_bytes": contact_sheet.stat().st_size,
            "sha256": file_sha256(contact_sheet),
        }

    return manifest


def write_manifest(image_files: list[Path]) -> Path:
    manifest = build_manifest(image_files)
    MANIFEST_OUTPUT.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return MANIFEST_OUTPUT


def main() -> None:
    image_files = get_image_files()

    print(f"Found {len(image_files)} image files, building PPTX...")

    prs = create_ppt(image_files)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPT_OUTPUT))

    manifest_path = write_manifest(image_files)

    print()
    print(f"PPTX created: {PPT_OUTPUT}")
    print(f"Slides: {len(image_files)}")
    print(f"Manifest: {manifest_path}")


if __name__ == "__main__":
    main()
